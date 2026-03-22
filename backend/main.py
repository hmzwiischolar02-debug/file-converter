"""
FileConvert — FastAPI Backend
=================================
All 7 converters use pure-Python libraries only.
No LibreOffice, no poppler, no system packages needed.

Install:
    pip install -r requirements.txt

Run:
    uvicorn main:app --reload --port 8000

Docs:
    http://localhost:8000/docs
"""

import io
import os
import time
import uuid
import zipfile
import asyncio
import threading
import tempfile
from pathlib import Path
from typing import List

from fastapi import FastAPI, File, Form, UploadFile, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse
from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.util import get_remote_address
from slowapi.errors import RateLimitExceeded
from slowapi.middleware import SlowAPIMiddleware

# ── Config ─────────────────────────────────────────────────────────────────────
UPLOAD_DIR = Path(tempfile.gettempdir()) / "fileconvert"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

MAX_FILE_SIZE   = 10 * 1024 * 1024   # 10 MB
FILE_TTL        = 3600                # seconds before auto-delete
# In production set CORS_ORIGIN env var to your Vercel URL
# e.g. CORS_ORIGIN=https://fileconvert.vercel.app
ALLOWED_ORIGINS = os.getenv("CORS_ORIGIN", "*").split(",")

# ── App ────────────────────────────────────────────────────────────────────────
app = FastAPI(
    title="FileConvert API",
    description="Free online file converter — PDF, Word, Images",
    version="2.0.0",
)

# ── Rate limiting ──────────────────────────────────────────────────────────────
# Free tier limits: 10 conversions per minute per IP
# Adjust via RATE_LIMIT env var e.g. "20/minute"
RATE_LIMIT = os.getenv("RATE_LIMIT", "10/minute")

limiter = Limiter(key_func=get_remote_address, default_limits=[RATE_LIMIT])
app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)
app.add_middleware(SlowAPIMiddleware)

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── Allowed extensions per tool ────────────────────────────────────────────────
ALLOWED = {
    "pdf-to-word":  {".pdf"},
    "word-to-pdf":  {".doc", ".docx"},
    "jpg-to-pdf":   {".jpg", ".jpeg", ".png"},
    "pdf-to-jpg":   {".pdf"},
    "png-to-ico":   {".png"},
    "merge-pdf":    {".pdf"},
    "compress-pdf": {".pdf"},
}

# ── Helpers ────────────────────────────────────────────────────────────────────

def validate(tool_id: str, filename: str) -> None:
    ext = Path(filename or "").suffix.lower()
    if ext not in ALLOWED.get(tool_id, set()):
        allowed_str = ", ".join(ALLOWED[tool_id])
        raise HTTPException(
            422,
            f"Invalid file type '{ext}'. Accepted: {allowed_str}",
        )


def validate_size(data: bytes) -> None:
    if len(data) > MAX_FILE_SIZE:
        raise HTTPException(
            413,
            f"File too large. Maximum size is {MAX_FILE_SIZE // 1024 // 1024} MB.",
        )


def save_tmp(data: bytes, suffix: str) -> Path:
    path = UPLOAD_DIR / f"{uuid.uuid4()}{suffix}"
    path.write_bytes(data)
    return path


def auto_delete(*paths: Path, delay: int = FILE_TTL) -> None:
    """Delete files after `delay` seconds in a background thread."""
    def _delete():
        time.sleep(delay)
        for p in paths:
            try:
                p.unlink(missing_ok=True)
            except Exception:
                pass
    threading.Thread(target=_delete, daemon=True).start()


def stream(data: bytes, filename: str, media_type: str) -> StreamingResponse:
    return StreamingResponse(
        io.BytesIO(data),
        media_type=media_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ── Conversion functions ───────────────────────────────────────────────────────

def do_pdf_to_word(pdf_bytes: bytes) -> bytes:
    """
    PDF → DOCX
    Primary:  pdf2docx  (best layout fidelity)
    Fallback: pypdf text extraction + python-docx
    """
    try:
        from pdf2docx import Converter as PDFConverter
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as src:
            src.write(pdf_bytes)
            src_path = src.name
        dst_path = src_path.replace(".pdf", ".docx")
        try:
            cv = PDFConverter(src_path)
            cv.convert(dst_path, start=0, end=None)
            cv.close()
            result = Path(dst_path).read_bytes()
        finally:
            Path(src_path).unlink(missing_ok=True)
            Path(dst_path).unlink(missing_ok=True)
        return result

    except ImportError:
        # Fallback — plain text extraction via pypdf + python-docx
        from pypdf import PdfReader
        from docx import Document
        from docx.shared import Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        reader = PdfReader(io.BytesIO(pdf_bytes))
        doc = Document()

        # Style heading
        doc.add_heading("Converted from PDF", level=1)

        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                # Add a page label paragraph
                label = doc.add_paragraph()
                label.add_run(f"— Page {i} —").bold = True
                label.alignment = WD_ALIGN_PARAGRAPH.CENTER

                for line in text.split("\n"):
                    if line.strip():
                        p = doc.add_paragraph(line.strip())
                        p.runs[0].font.size = Pt(11)

            if i < len(reader.pages):
                doc.add_page_break()

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()


def do_word_to_pdf(docx_bytes: bytes, original_name: str) -> bytes:
    """
    DOCX/DOC → PDF — conversion chain:

    1. docx2pdf      → MS Word COM on Windows (perfect, if Word is installed)
    2. LibreOffice   → perfect on Linux/macOS
    3. mammoth + xhtml2pdf → pure Python, works on ALL platforms (good quality)
    4. python-docx + reportlab → basic text fallback (always works)
    """
    import logging, subprocess
    logger = logging.getLogger("fileconvert")
    suffix = Path(original_name).suffix.lower()

    # ── 1. docx2pdf — only works on Windows/macOS with MS Word installed ────────
    # Skipped on Linux (Railway/Render) where MS Word is not available
    import platform
    if platform.system() in ("Windows", "Darwin"):
        try:
            from docx2pdf import convert as d2p
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                tmp.write(docx_bytes)
                src = tmp.name
            dst = src.replace(suffix, ".pdf")
            try:
                d2p(src, dst)
                if Path(dst).exists() and Path(dst).stat().st_size > 100:
                    logger.info("word-to-pdf: docx2pdf (MS Word)")
                    return Path(dst).read_bytes()
            finally:
                Path(src).unlink(missing_ok=True)
                Path(dst).unlink(missing_ok=True)
        except ImportError:
            pass
        except Exception as e:
            logger.warning(f"word-to-pdf: docx2pdf failed — {e}")

    # ── 2. LibreOffice CLI (Linux / macOS) ────────────────────────────────────
    try:
        import shutil
        if shutil.which("libreoffice") or shutil.which("soffice"):
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                tmp.write(docx_bytes)
                src = tmp.name
            out_dir = tempfile.mkdtemp()
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf",
                 "--outdir", out_dir, src],
                capture_output=True, timeout=60,
            )
            pdf_path = Path(out_dir) / (Path(src).stem + ".pdf")
            if pdf_path.exists() and pdf_path.stat().st_size > 100:
                data = pdf_path.read_bytes()
                Path(src).unlink(missing_ok=True)
                pdf_path.unlink(missing_ok=True)
                logger.info("word-to-pdf: LibreOffice")
                return data
    except Exception as e:
        logger.warning(f"word-to-pdf: LibreOffice failed — {e}")

    # ── 3. mammoth → HTML → xhtml2pdf (pure Python, works on Windows) ─────────
    try:
        import mammoth
        from xhtml2pdf import pisa

        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
            tmp.write(docx_bytes)
            src = tmp.name

        try:
            with open(src, "rb") as f:
                result = mammoth.convert_to_html(f)

            html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8"/>
<style>
  @page {{
    size: A4;
    margin: 2.5cm;
  }}
  body {{
    font-family: Helvetica, Arial, sans-serif;
    font-size: 12pt;
    line-height: 1.6;
    color: #000;
  }}
  h1 {{ font-size: 20pt; font-weight: bold; margin: 14pt 0 8pt 0; color: #1e293b; }}
  h2 {{ font-size: 16pt; font-weight: bold; margin: 12pt 0 6pt 0; color: #1e293b; }}
  h3 {{ font-size: 13pt; font-weight: bold; margin: 10pt 0 4pt 0; color: #334155; }}
  p  {{ margin: 0 0 8pt 0; }}
  table {{ border-collapse: collapse; width: 100%; margin: 10pt 0; }}
  td, th {{ border: 1px solid #cbd5e1; padding: 5pt 7pt; font-size: 10pt; }}
  th {{ background-color: #2563eb; color: white; font-weight: bold; }}
  ul, ol {{ margin: 0 0 8pt 18pt; padding: 0; }}
  li {{ margin-bottom: 3pt; }}
  strong, b {{ font-weight: bold; }}
  em, i {{ font-style: italic; }}
  u {{ text-decoration: underline; }}
  img {{ max-width: 100%; }}
</style>
</head>
<body>{result.value}</body>
</html>"""

            buf = io.BytesIO()
            pisa_status = pisa.CreatePDF(html, dest=buf, encoding="utf-8")
            if not pisa_status.err:
                logger.info("word-to-pdf: mammoth + xhtml2pdf")
                return buf.getvalue()
            else:
                logger.warning(f"word-to-pdf: xhtml2pdf error code {pisa_status.err}")

        finally:
            Path(src).unlink(missing_ok=True)

    except ImportError as e:
        logger.warning(f"word-to-pdf: mammoth/xhtml2pdf not installed — {e}")
    except Exception as e:
        logger.warning(f"word-to-pdf: mammoth/xhtml2pdf failed — {e}")

    # ── 4. python-docx + reportlab (always works, basic quality) ──────────────
    try:
        from docx import Document as DocxDoc
        from docx.oxml.ns import qn
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_CENTER, TA_RIGHT, TA_JUSTIFY
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        )

        def esc(t):
            return (t.replace("&", "&amp;")
                     .replace("<", "&lt;")
                     .replace(">", "&gt;")
                     .replace('"', "&quot;"))

        def runs_to_rich(para):
            parts = []
            for run in para.runs:
                txt = esc(run.text)
                if not txt:
                    continue
                o, c = "", ""
                try:
                    if run.font.size and run.font.size.pt:
                        sz = int(run.font.size.pt)
                        o += f'<font size="{sz}">'; c = "</font>" + c
                except Exception:
                    pass
                try:
                    if run.font.color and run.font.color.type is not None:
                        rgb = str(run.font.color.rgb)
                        o += f'<font color="#{rgb}">'; c = "</font>" + c
                except Exception:
                    pass
                if run.bold:      o += "<b>";  c = "</b>"  + c
                if run.italic:    o += "<i>";  c = "</i>"  + c
                if run.underline: o += "<u>";  c = "</u>"  + c
                parts.append(f"{o}{txt}{c}")
            return "".join(parts)

        doc  = DocxDoc(io.BytesIO(docx_bytes))
        buf  = io.BytesIO()
        pdf  = SimpleDocTemplate(
            buf, pagesize=A4,
            leftMargin=2.5*cm, rightMargin=2.5*cm,
            topMargin=2.5*cm,  bottomMargin=2.5*cm,
        )
        base = getSampleStyleSheet()
        ST = {
            "h1":   ParagraphStyle("H1",  parent=base["Normal"], fontSize=20,
                                   fontName="Helvetica-Bold", spaceAfter=10, spaceBefore=14),
            "h2":   ParagraphStyle("H2",  parent=base["Normal"], fontSize=16,
                                   fontName="Helvetica-Bold", spaceAfter=8,  spaceBefore=10),
            "h3":   ParagraphStyle("H3",  parent=base["Normal"], fontSize=13,
                                   fontName="Helvetica-Bold", spaceAfter=6,  spaceBefore=8),
            "body": ParagraphStyle("Body",parent=base["Normal"], fontSize=11,
                                   leading=17, spaceAfter=5),
            "bull": ParagraphStyle("Bull",parent=base["Normal"], fontSize=11,
                                   leading=17, leftIndent=18, spaceAfter=3),
        }
        ALIGN = {"CENTER": TA_CENTER, "RIGHT": TA_RIGHT, "JUSTIFY": TA_JUSTIFY}
        story = []

        for para in doc.paragraphs:
            sn  = (para.style.name or "").strip()
            raw = para.text
            if not raw.strip():
                story.append(Spacer(1, 5)); continue
            if   "Heading 1" in sn or sn == "Title": st = ST["h1"]
            elif "Heading 2" in sn:                  st = ST["h2"]
            elif "Heading 3" in sn:                  st = ST["h3"]
            else:                                     st = ST["body"]
            al = para.alignment.name if para.alignment else None
            if al in ALIGN:
                st = ParagraphStyle(f"{st.name}_a", parent=st, alignment=ALIGN[al])
            rich = runs_to_rich(para) or esc(raw)
            is_bullet = ("List" in sn or
                         para._element.find(qn("w:numPr")) is not None)
            story.append(
                Paragraph(f"• &nbsp;{rich}", ST["bull"])
                if is_bullet else Paragraph(rich, st)
            )

        for table in doc.tables:
            td = []
            for row in table.rows:
                td.append([Paragraph(esc(c.text.strip()), ST["body"]) for c in row.cells])
            if td:
                cw = (A4[0] - 5*cm) / max(len(r) for r in td)
                t  = Table(td, colWidths=[cw] * max(len(r) for r in td), repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND",    (0,0),(-1,0), colors.HexColor("#2563eb")),
                    ("TEXTCOLOR",     (0,0),(-1,0), colors.white),
                    ("FONTNAME",      (0,0),(-1,0), "Helvetica-Bold"),
                    ("FONTSIZE",      (0,0),(-1,-1), 9),
                    ("GRID",          (0,0),(-1,-1), 0.4, colors.HexColor("#cbd5e1")),
                    ("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white, colors.HexColor("#f8fafc")]),
                    ("TOPPADDING",    (0,0),(-1,-1), 5),
                    ("BOTTOMPADDING", (0,0),(-1,-1), 5),
                    ("LEFTPADDING",   (0,0),(-1,-1), 6),
                    ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
                ]))
                story.append(Spacer(1, 10))
                story.append(t)
                story.append(Spacer(1, 10))

        if not story:
            story.append(Paragraph("(empty document)", ST["body"]))

        pdf.build(story)
        logger.info("word-to-pdf: reportlab fallback")
        return buf.getvalue()

    except Exception as e:
        raise HTTPException(500, f"Word to PDF failed at all stages. Last error: {e}")



def do_images_to_pdf(images: List[bytes], names: List[str]) -> bytes:
    """
    JPG / PNG (multiple) → PDF
    Uses Pillow + reportlab — no external tools needed.
    """
    from PIL import Image
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.platypus import SimpleDocTemplate, Image as RLImage, PageBreak
    import tempfile, os

    buf = io.BytesIO()
    story = []
    tmp_files = []

    page_w, page_h = A4
    margin = 20

    for img_bytes, name in zip(images, names):
        with Image.open(io.BytesIO(img_bytes)) as img:
            # Convert to RGB (handles RGBA, P, L modes)
            if img.mode != "RGB":
                background = Image.new("RGB", img.size, (255, 255, 255))
                if img.mode in ("RGBA", "LA"):
                    background.paste(img, mask=img.split()[-1])
                else:
                    background.paste(img.convert("RGBA"), mask=None)
                img = background

            img_w, img_h = img.size

            # Scale to fit page
            max_w = page_w - 2 * margin
            max_h = page_h - 2 * margin
            ratio = min(max_w / img_w, max_h / img_h, 1.0)
            new_w = img_w * ratio
            new_h = img_h * ratio

            # Save to temp file so reportlab can read it
            tmp = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
            img.save(tmp.name, "JPEG", quality=90, optimize=True)
            tmp.close()
            tmp_files.append(tmp.name)

            story.append(RLImage(tmp.name, width=new_w, height=new_h))
            story.append(PageBreak())

    # Remove trailing page break
    if story and isinstance(story[-1], PageBreak):
        story.pop()

    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=margin, rightMargin=margin,
        topMargin=margin,  bottomMargin=margin,
    )
    doc.build(story)

    # Clean up temp files
    for f in tmp_files:
        try:
            os.unlink(f)
        except Exception:
            pass

    return buf.getvalue()


def do_pdf_to_images(pdf_bytes: bytes) -> bytes:
    """
    PDF → JPG images (one per page) → ZIP
    Primary:  PyMuPDF (fitz) — pure Python wheel, no poppler needed
    Fallback: pdf2image (requires poppler)
    Fallback: pypdf embedded image extraction
    """
    zip_buf = io.BytesIO()

    # ── Primary: PyMuPDF ──────────────────────────────────────────────────────
    try:
        import fitz  # PyMuPDF

        pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for page_num in range(len(pdf_doc)):
                page = pdf_doc[page_num]
                # 2x scale for 150 DPI equivalent
                mat = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img_bytes = pix.tobytes("jpeg")
                zf.writestr(f"page_{page_num + 1}.jpg", img_bytes)
        pdf_doc.close()
        return zip_buf.getvalue()

    except ImportError:
        pass

    # ── Fallback: pdf2image (needs poppler) ───────────────────────────────────
    try:
        from pdf2image import convert_from_bytes
        images = convert_from_bytes(pdf_bytes, dpi=150)
        with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, img in enumerate(images):
                buf = io.BytesIO()
                img.save(buf, "JPEG", quality=90)
                zf.writestr(f"page_{i + 1}.jpg", buf.getvalue())
        return zip_buf.getvalue()

    except (ImportError, Exception):
        pass

    # ── Last resort: extract embedded images via pypdf ─────────────────────────
    from pypdf import PdfReader
    from PIL import Image as PILImage

    reader = PdfReader(io.BytesIO(pdf_bytes))
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        img_count = 0
        for page_num, page in enumerate(reader.pages):
            for img_obj in page.images:
                try:
                    pil = PILImage.open(io.BytesIO(img_obj.data))
                    if pil.mode != "RGB":
                        pil = pil.convert("RGB")
                    buf = io.BytesIO()
                    pil.save(buf, "JPEG", quality=90)
                    img_count += 1
                    zf.writestr(f"page_{page_num + 1}_img_{img_count}.jpg", buf.getvalue())
                except Exception:
                    continue

        if img_count == 0:
            raise HTTPException(
                422,
                "Could not extract images from this PDF. "
                "Install PyMuPDF for full page rendering: pip install PyMuPDF",
            )

    return zip_buf.getvalue()


def do_png_to_ico(png_bytes: bytes) -> bytes:
    """PNG → ICO with multiple sizes (pure Pillow)."""
    from PIL import Image

    with Image.open(io.BytesIO(png_bytes)) as img:
        img = img.convert("RGBA")
        buf = io.BytesIO()
        img.save(
            buf,
            format="ICO",
            sizes=[(16, 16), (32, 32), (48, 48), (64, 64), (128, 128)],
        )
        return buf.getvalue()


def do_merge_pdfs(pdf_list: List[bytes]) -> bytes:
    """Merge multiple PDFs into one (pure pypdf)."""
    from pypdf import PdfWriter, PdfReader

    writer = PdfWriter()
    for pdf_bytes in pdf_list:
        reader = PdfReader(io.BytesIO(pdf_bytes))
        for page in reader.pages:
            writer.add_page(page)

    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def do_compress_pdf(pdf_bytes: bytes) -> bytes:
    """
    Compress PDF — reduces size via content stream compression + deduplication.
    Pure pypdf, no external tools.

    Key fix: pages must be added to the writer FIRST, then compressed —
    compress_content_streams() requires the page to already belong to a PdfWriter.
    """
    from pypdf import PdfWriter, PdfReader

    reader = PdfReader(io.BytesIO(pdf_bytes))
    writer = PdfWriter()

    # Step 1 — add all pages to the writer
    for page in reader.pages:
        writer.add_page(page)

    # Step 2 — compress content streams on the writer's own pages
    for page in writer.pages:
        page.compress_content_streams()

    # Step 3 — deduplicate shared objects (images, fonts)
    try:
        writer.compress_identical_objects(
            remove_identicals=True,
            remove_orphans=True,
        )
    except Exception:
        pass  # non-critical optimisation; skip if unsupported version

    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def do_excel_to_pdf(xlsx_bytes: bytes, filename: str) -> bytes:
    """
    Excel (.xls / .xlsx) → PDF
    Primary:  openpyxl + reportlab (pure Python, all platforms)
    """
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (
        SimpleDocTemplate, Table, TableStyle, Paragraph,
        Spacer, PageBreak,
    )
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

    wb = load_workbook(io.BytesIO(xlsx_bytes), data_only=True)
    buf = io.BytesIO()
    styles_rl = getSampleStyleSheet()
    cell_style = ParagraphStyle("Cell", parent=styles_rl["Normal"],
                                fontSize=8, leading=11)
    header_style = ParagraphStyle("Header", parent=styles_rl["Normal"],
                                  fontSize=8, fontName="Helvetica-Bold", leading=11)
    story = []

    for sheet_idx, ws in enumerate(wb.worksheets):
        # Sheet title
        title_style = ParagraphStyle("Title", parent=styles_rl["Normal"],
                                     fontSize=12, fontName="Helvetica-Bold",
                                     spaceAfter=8)
        story.append(Paragraph(ws.title, title_style))

        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            story.append(Paragraph("(empty sheet)", cell_style))
            if sheet_idx < len(wb.worksheets) - 1:
                story.append(PageBreak())
            continue

        # Filter out fully empty rows
        rows = [r for r in rows if any(c is not None for c in r)]
        if not rows:
            continue

        # Build table data
        data = []
        for row_idx, row in enumerate(rows):
            row_data = []
            for cell in row:
                txt = str(cell) if cell is not None else ""
                st = header_style if row_idx == 0 else cell_style
                row_data.append(Paragraph(txt[:200], st))  # cap at 200 chars
            data.append(row_data)

        if not data:
            continue

        col_count = max(len(r) for r in data)
        page_w = landscape(A4)[0] - 2 * cm
        col_w = page_w / col_count if col_count else page_w

        t = Table(data, colWidths=[col_w] * col_count, repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND",    (0, 0), (-1, 0),  colors.HexColor("#16a34a")),
            ("TEXTCOLOR",     (0, 0), (-1, 0),  colors.white),
            ("FONTNAME",      (0, 0), (-1, 0),  "Helvetica-Bold"),
            ("FONTSIZE",      (0, 0), (-1, -1), 8),
            ("GRID",          (0, 0), (-1, -1), 0.3, colors.HexColor("#d1d5db")),
            ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, colors.HexColor("#f9fafb")]),
            ("TOPPADDING",    (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ("LEFTPADDING",   (0, 0), (-1, -1), 4),
            ("VALIGN",        (0, 0), (-1, -1), "TOP"),
        ]))

        story.append(t)
        if sheet_idx < len(wb.worksheets) - 1:
            story.append(PageBreak())
        else:
            story.append(Spacer(1, 8))

    doc = SimpleDocTemplate(buf, pagesize=landscape(A4),
                            leftMargin=cm, rightMargin=cm,
                            topMargin=cm, bottomMargin=cm)
    doc.build(story)
    return buf.getvalue()


def do_ppt_to_pdf(pptx_bytes: bytes) -> bytes:
    """
    PowerPoint (.ppt / .pptx) → PDF
    Primary:  python-pptx → extract text per slide + reportlab
    """
    from pptx import Presentation
    from pptx.util import Pt
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, HRFlowable
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

    prs = Presentation(io.BytesIO(pptx_bytes))
    buf = io.BytesIO()

    base = getSampleStyleSheet()
    slide_num_style = ParagraphStyle("SlideNum", parent=base["Normal"],
                                     fontSize=8, textColor=colors.HexColor("#94a3b8"),
                                     spaceAfter=4)
    title_style = ParagraphStyle("SlideTitle", parent=base["Normal"],
                                  fontSize=18, fontName="Helvetica-Bold",
                                  textColor=colors.HexColor("#1e293b"),
                                  spaceAfter=8, spaceBefore=4)
    body_style = ParagraphStyle("SlideBody", parent=base["Normal"],
                                 fontSize=11, leading=16,
                                 textColor=colors.HexColor("#334155"),
                                 spaceAfter=4)
    bullet_style = ParagraphStyle("SlideBullet", parent=base["Normal"],
                                   fontSize=10, leading=15,
                                   textColor=colors.HexColor("#475569"),
                                   leftIndent=16, spaceAfter=3)

    def esc(t):
        return (t.replace("&","&amp;")
                 .replace("<","&lt;")
                 .replace(">","&gt;"))

    story = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        story.append(Paragraph(f"Slide {slide_idx}", slide_num_style))

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if not text:
                    continue
                safe = esc(text)
                # Heuristic: first text box first paragraph = title
                if para_idx == 0 and shape == slide.shapes[0]:
                    story.append(Paragraph(safe, title_style))
                elif para.level > 0:
                    indent = "  " * para.level
                    story.append(Paragraph(f"• {safe}", bullet_style))
                else:
                    story.append(Paragraph(safe, body_style))

        story.append(HRFlowable(width="100%", thickness=0.5,
                                color=colors.HexColor("#e2e8f0"),
                                spaceAfter=6, spaceBefore=6))
        if slide_idx < len(prs.slides):
            story.append(PageBreak())

    if not story:
        story.append(Paragraph("(empty presentation)", body_style))

    doc = SimpleDocTemplate(buf, pagesize=landscape(A4),
                            leftMargin=2*cm, rightMargin=2*cm,
                            topMargin=1.5*cm, bottomMargin=1.5*cm)
    doc.build(story)
    return buf.getvalue()


def do_pdf_to_png(pdf_bytes: bytes) -> bytes:
    """PDF → PNG images (one per page) → ZIP  — uses PyMuPDF."""
    import fitz
    import zipfile

    pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    zip_buf = io.BytesIO()

    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            mat  = fitz.Matrix(2.0, 2.0)        # 150 DPI equivalent
            pix  = page.get_pixmap(matrix=mat, alpha=False)
            zf.writestr(f"page_{page_num + 1}.png", pix.tobytes("png"))

    pdf_doc.close()
    return zip_buf.getvalue()


def do_compress_image(img_bytes: bytes, filename: str) -> tuple:
    """
    Compress JPG / PNG image.
    Returns (compressed_bytes, output_filename).
    """
    from PIL import Image

    ext = Path(filename).suffix.lower()

    with Image.open(io.BytesIO(img_bytes)) as img:
        original_size = len(img_bytes)
        out_buf = io.BytesIO()

        if ext in (".jpg", ".jpeg"):
            # Convert to RGB if needed (RGBA not supported in JPEG)
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            img.save(out_buf, format="JPEG", quality=75, optimize=True)
            out_name = Path(filename).stem + "_compressed.jpg"
        else:
            # PNG — use maximum compression
            if img.mode == "P":
                img = img.convert("RGBA")
            img.save(out_buf, format="PNG", optimize=True, compress_level=9)
            out_name = Path(filename).stem + "_compressed.png"

        return out_buf.getvalue(), out_name


def do_unlock_pdf(pdf_bytes: bytes, password: str) -> bytes:
    """Remove password protection from a PDF using pypdf."""
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(io.BytesIO(pdf_bytes))

    if reader.is_encrypted:
        success = reader.decrypt(password)
        if not success:
            raise HTTPException(400, "Wrong password. Please check and try again.")

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


# ── Routes ─────────────────────────────────────────────────────────────────────

def do_jpg_to_png(img_bytes: bytes) -> bytes:
    """JPG → PNG (lossless, preserves full quality)."""
    from PIL import Image
    with Image.open(io.BytesIO(img_bytes)) as img:
        buf = io.BytesIO()
        img.convert("RGBA").save(buf, format="PNG", optimize=True)
        return buf.getvalue()


def do_png_to_jpg(img_bytes: bytes) -> bytes:
    """PNG → JPG (flatten transparency to white, 90% quality)."""
    from PIL import Image
    with Image.open(io.BytesIO(img_bytes)) as img:
        # Flatten transparent background to white
        background = Image.new("RGB", img.size, (255, 255, 255))
        if img.mode in ("RGBA", "LA", "P"):
            background.paste(img.convert("RGBA"), mask=img.convert("RGBA").split()[3])
        else:
            background.paste(img.convert("RGB"))
        buf = io.BytesIO()
        background.save(buf, format="JPEG", quality=90, optimize=True)
        return buf.getvalue()


def do_html_to_pdf(html_bytes: bytes) -> bytes:
    """
    HTML → PDF with full CSS support.

    Chain:
    1. WeasyPrint  — full CSS, perfect rendering (works on Linux/Railway)
    2. xhtml2pdf   — good CSS support, pure Python fallback
    3. reportlab   — basic text extraction, last resort
    """
    html_content = html_bytes.decode("utf-8", errors="replace")

    # Ensure complete HTML structure with good base styles
    if "<html" not in html_content.lower():
        html_content = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8"/>
<style>
  @page {{ margin: 2cm; }}
  body {{
    font-family: Arial, Helvetica, sans-serif;
    font-size: 12pt;
    line-height: 1.6;
    color: #000;
  }}
  h1 {{ font-size: 22pt; font-weight: bold; margin: 16pt 0 8pt; color: #1e293b; }}
  h2 {{ font-size: 18pt; font-weight: bold; margin: 14pt 0 6pt; color: #1e293b; }}
  h3 {{ font-size: 14pt; font-weight: bold; margin: 12pt 0 4pt; color: #334155; }}
  p  {{ margin: 0 0 8pt; }}
  table {{ border-collapse: collapse; width: 100%; margin: 10pt 0; }}
  td, th {{ border: 1px solid #cbd5e1; padding: 6pt 8pt; font-size: 10pt; }}
  th {{ background-color: #2563eb; color: white; font-weight: bold; }}
  ul, ol {{ margin: 0 0 8pt 20pt; }}
  li {{ margin-bottom: 4pt; }}
  strong, b {{ font-weight: bold; }}
  em, i {{ font-style: italic; }}
  u {{ text-decoration: underline; }}
  code, pre {{ font-family: monospace; background: #f8fafc;
               padding: 2pt 4pt; border-radius: 3pt; font-size: 10pt; }}
  blockquote {{ border-left: 3pt solid #2563eb; margin: 8pt 0;
                padding: 4pt 12pt; color: #475569; }}
  img {{ max-width: 100%; }}
  a {{ color: #2563eb; }}
</style>
</head>
<body>{html_content}</body>
</html>"""

    # ── 1. WeasyPrint — best CSS support, works on Linux ──────────────────────
    try:
        from weasyprint import HTML as WeasyHTML, CSS
        from weasyprint.text.fonts import FontConfiguration

        font_config = FontConfiguration()
        pdf_bytes = WeasyHTML(
            string=html_content,
            base_url=None,
        ).write_pdf(font_config=font_config)

        if pdf_bytes and len(pdf_bytes) > 500:
            return pdf_bytes

    except ImportError:
        pass  # weasyprint not installed
    except Exception as e:
        pass  # fall through

    # ── 2. xhtml2pdf fallback ─────────────────────────────────────────────────
    try:
        from xhtml2pdf import pisa
        buf = io.BytesIO()
        status = pisa.CreatePDF(html_content, dest=buf, encoding="utf-8")
        if not status.err and len(buf.getvalue()) > 500:
            return buf.getvalue()
    except Exception:
        pass

    # ── 3. reportlab plain text — last resort ─────────────────────────────────
    import re
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    text = re.sub(r"<[^>]+>", " ", html_content)
    for ent, rep in [("&nbsp;"," "),("&amp;","&"),("&lt;","<"),("&gt;",">"),("&quot;",'"')]:
        text = text.replace(ent, rep)
    text = re.sub(r"\s+", " ", text).strip()

    buf = io.BytesIO()
    styles = getSampleStyleSheet()
    body = ParagraphStyle("B", parent=styles["Normal"], fontSize=11, leading=17)
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=2.5*cm, rightMargin=2.5*cm,
                            topMargin=2.5*cm,  bottomMargin=2.5*cm)
    story = []
    for line in text.split(". "):
        line = line.strip()
        if line:
            safe = line.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
            story.append(Paragraph(safe, body))
            story.append(Spacer(1, 4))
    if not story:
        story.append(Paragraph("(empty document)", body))
    doc.build(story)
    return buf.getvalue()


def do_watermark_pdf(pdf_bytes: bytes, watermark_text: str) -> bytes:
    """Add diagonal text watermark to every page of a PDF."""
    from pypdf import PdfWriter, PdfReader
    from reportlab.lib.pagesizes import letter
    from reportlab.lib import colors
    from reportlab.pdfgen import canvas as rl_canvas
    import math

    # Create a single-page watermark PDF in memory
    wm_buf = io.BytesIO()
    c = rl_canvas.Canvas(wm_buf, pagesize=letter)
    page_w, page_h = letter

    c.saveState()
    c.setFont("Helvetica-Bold", 48)
    c.setFillColor(colors.Color(0.6, 0.6, 0.6, alpha=0.3))  # 30% opacity grey
    c.translate(page_w / 2, page_h / 2)
    c.rotate(45)
    c.drawCentredString(0, 0, watermark_text.upper())
    c.restoreState()
    c.save()
    wm_buf.seek(0)

    from pypdf import PdfReader as PR
    wm_reader = PR(wm_buf)
    wm_page   = wm_reader.pages[0]

    reader = PdfReader(io.BytesIO(pdf_bytes))
    writer = PdfWriter()

    for page in reader.pages:
        # Merge watermark onto each page
        page.merge_page(wm_page)
        writer.add_page(page)

    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def do_pdf_to_html(pdf_bytes: bytes) -> bytes:
    """
    PDF → HTML using PyMuPDF for best text extraction with layout.
    Returns HTML bytes.
    """
    import fitz

    pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    pages_html = []

    for page_num in range(len(pdf_doc)):
        page = pdf_doc[page_num]
        # get_text("html") gives structured HTML per page
        page_html = page.get_text("html")
        pages_html.append(f"""
        <div class="page" id="page-{page_num + 1}">
          <div class="page-label">Page {page_num + 1}</div>
          {page_html}
        </div>
        """)

    pdf_doc.close()

    full_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8"/>
  <meta name="viewport" content="width=device-width, initial-scale=1.0"/>
  <title>Converted PDF</title>
  <style>
    body {{
      font-family: Arial, Helvetica, sans-serif;
      font-size: 12pt;
      line-height: 1.6;
      color: #000;
      max-width: 860px;
      margin: 0 auto;
      padding: 20px;
      background: #f8fafc;
    }}
    .page {{
      background: #fff;
      border: 1px solid #e2e8f0;
      border-radius: 8px;
      padding: 40px 48px;
      margin-bottom: 32px;
      box-shadow: 0 2px 8px rgba(0,0,0,0.06);
    }}
    .page-label {{
      font-size: 10pt;
      color: #94a3b8;
      font-weight: 600;
      margin-bottom: 16px;
      padding-bottom: 8px;
      border-bottom: 1px solid #f1f5f9;
      text-transform: uppercase;
      letter-spacing: .05em;
    }}
    p {{ margin: 0 0 8pt; }}
    span {{ font-family: inherit !important; }}
  </style>
</head>
<body>
  {"".join(pages_html)}
  <footer style="text-align:center;padding:20px;color:#94a3b8;font-size:10pt;">
    Converted with FileConvert.io
  </footer>
</body>
</html>"""

    return full_html.encode("utf-8")


def do_remove_background(img_bytes: bytes) -> bytes:
    """
    Remove image background.
    Primary:  rembg (AI, optional)
    Fallback: Pillow GrabCut-style with size limit to prevent hanging
    """
    from PIL import Image, ImageFilter

    # ── Primary: rembg (if installed) ────────────────────────────────────────
    try:
        import importlib
        rembg_mod = importlib.import_module("rembg")
        with Image.open(io.BytesIO(img_bytes)) as img:
            output = rembg_mod.remove(img.convert("RGBA"))
            buf = io.BytesIO()
            output.save(buf, format="PNG")
            return buf.getvalue()
    except (ImportError, ModuleNotFoundError):
        pass
    except Exception:
        pass

    # ── Fallback: Pillow flood-fill (capped at 800px to prevent hanging) ─────
    with Image.open(io.BytesIO(img_bytes)) as img:
        img = img.convert("RGBA")
        w, h = img.size

        # Resize large images to max 800px to prevent BFS from hanging
        MAX_DIM = 800
        scale = 1.0
        if w > MAX_DIM or h > MAX_DIM:
            scale = MAX_DIM / max(w, h)
            new_w = int(w * scale)
            new_h = int(h * scale)
            img = img.resize((new_w, new_h), Image.LANCZOS)
            w, h = new_w, new_h

        pixels = img.load()

        # Sample background colour from 4 corners
        corners = [
            pixels[0, 0], pixels[w-1, 0],
            pixels[0, h-1], pixels[w-1, h-1],
        ]
        bg_r = sum(c[0] for c in corners) // 4
        bg_g = sum(c[1] for c in corners) // 4
        bg_b = sum(c[2] for c in corners) // 4
        tolerance = 35

        def is_bg(r, g, b, a):
            return (abs(r - bg_r) < tolerance and
                    abs(g - bg_g) < tolerance and
                    abs(b - bg_b) < tolerance)

        # BFS flood fill from all edges — capped at 500k pixels max
        from collections import deque
        visited = set()
        queue = deque()
        MAX_PIXELS = 500_000

        # Seed from all 4 edges
        for x in range(w):
            for y in [0, h-1]:
                if (x, y) not in visited:
                    r, g, b, a = pixels[x, y]
                    if is_bg(r, g, b, a):
                        queue.append((x, y))
                        visited.add((x, y))
        for y in range(h):
            for x in [0, w-1]:
                if (x, y) not in visited:
                    r, g, b, a = pixels[x, y]
                    if is_bg(r, g, b, a):
                        queue.append((x, y))
                        visited.add((x, y))

        processed = 0
        while queue and processed < MAX_PIXELS:
            x, y = queue.popleft()
            r, g, b, a = pixels[x, y]
            if not is_bg(r, g, b, a):
                continue
            pixels[x, y] = (r, g, b, 0)  # transparent
            processed += 1
            for dx, dy in [(-1,0),(1,0),(0,-1),(0,1)]:
                nx, ny = x+dx, y+dy
                if 0 <= nx < w and 0 <= ny < h and (nx, ny) not in visited:
                    visited.add((nx, ny))
                    queue.append((nx, ny))

        # If we scaled down, scale back up
        if scale < 1.0:
            orig_w = int(w / scale)
            orig_h = int(h / scale)
            img = img.resize((orig_w, orig_h), Image.LANCZOS)

        img = img.filter(ImageFilter.SMOOTH)
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()


# ── Diagnostics endpoint ───────────────────────────────────────────────────────
@app.get("/diagnostics")
async def diagnostics():
    """
    Visit http://localhost:8000/diagnostics to see which conversion
    methods are available on your system.
    """
    import subprocess, shutil
    results = {}

    # docx2pdf (best for Windows)
    try:
        import docx2pdf
        results["docx2pdf"] = {"available": True, "note": "✅ BEST — uses MS Word directly"}
    except ImportError:
        results["docx2pdf"] = {"available": False, "note": "❌ Run: pip install docx2pdf"}

    # mammoth
    try:
        import mammoth
        results["mammoth"] = {"available": True, "note": "✅ Good HTML extraction"}
    except ImportError:
        results["mammoth"] = {"available": False, "note": "❌ Run: pip install mammoth"}

    # weasyprint
    try:
        from weasyprint import HTML
        results["weasyprint"] = {"available": True, "note": "✅ Good HTML→PDF"}
    except ImportError:
        results["weasyprint"] = {"available": False, "note": "❌ Run: pip install weasyprint"}

    # LibreOffice
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    results["libreoffice"] = {
        "available": bool(lo),
        "path": lo,
        "note": "✅ Excellent" if lo else "❌ Not installed (optional on Linux/macOS)"
    }

    # PyMuPDF (for PDF→JPG)
    try:
        import fitz
        results["PyMuPDF"] = {"available": True, "version": fitz.version[0], "note": "✅ Used for PDF→JPG"}
    except ImportError:
        results["PyMuPDF"] = {"available": False, "note": "❌ Run: pip install PyMuPDF"}

    # Determine active Word→PDF method (docx2pdf only works on Windows/macOS)
    import platform
    is_windows_or_mac = platform.system() in ("Windows", "Darwin")
    if is_windows_or_mac and results["docx2pdf"]["available"]:
        active = "docx2pdf (MS Word COM) — PERFECT quality"
    elif results["libreoffice"]["available"]:
        active = "LibreOffice — PERFECT quality"
    elif results["mammoth"]["available"]:
        active = "mammoth + xhtml2pdf — GOOD quality (active on this server)"
    else:
        active = "❌ No method available — install mammoth + xhtml2pdf"

    return {
        "word_to_pdf_active_method": active,
        "libraries": results,
        "recommendation": "pip install docx2pdf mammoth weasyprint PyMuPDF"
    }


@app.get("/health")
async def health():
    return {"status": "ok", "version": "2.0.0"}


@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    return JSONResponse(
        status_code=500,
        content={"detail": f"Unexpected error: {str(exc)}"},
    )


# ── 1. PDF → Word ──────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/pdf-to-word")
async def pdf_to_word(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    validate("pdf-to-word", f.filename)
    raw = await f.read()
    validate_size(raw)

    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_pdf_to_word, raw
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"PDF to Word conversion failed: {e}")

    return stream(
        data,
        "converted.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


# ── 2. Word → PDF ──────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/word-to-pdf")
async def word_to_pdf(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    fname = f.filename or "document.docx"
    validate("word-to-pdf", fname)
    raw = await f.read()
    validate_size(raw)

    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_word_to_pdf, raw, fname
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Word to PDF conversion failed: {e}")

    return stream(data, "converted.pdf", "application/pdf")


# ── 3. JPG/PNG → PDF ───────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/jpg-to-pdf")
async def jpg_to_pdf(request: Request, files: List[UploadFile] = File(...)):
    if not files:
        raise HTTPException(422, "No files provided.")

    image_data, image_names = [], []
    for f in files:
        validate("jpg-to-pdf", f.filename)
        raw = await f.read()
        validate_size(raw)
        image_data.append(raw)
        image_names.append(f.filename or "image.jpg")

    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_images_to_pdf, image_data, image_names
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Image to PDF conversion failed: {e}")

    return stream(data, "images.pdf", "application/pdf")


# ── 4. PDF → JPG ───────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/pdf-to-jpg")
async def pdf_to_jpg(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    validate("pdf-to-jpg", f.filename)
    raw = await f.read()
    validate_size(raw)

    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_pdf_to_images, raw
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"PDF to JPG conversion failed: {e}")

    return stream(data, "pages.zip", "application/zip")


# ── 5. PNG → ICO ───────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/png-to-ico")
async def png_to_ico(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    validate("png-to-ico", f.filename)
    raw = await f.read()
    validate_size(raw)

    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_png_to_ico, raw
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"PNG to ICO conversion failed: {e}")

    return stream(data, "favicon.ico", "image/x-icon")


# ── 6. Merge PDFs ──────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/merge-pdf")
async def merge_pdf(request: Request, files: List[UploadFile] = File(...)):
    if len(files) < 2:
        raise HTTPException(422, "Please upload at least 2 PDF files to merge.")

    pdf_list = []
    for f in files:
        validate("merge-pdf", f.filename)
        raw = await f.read()
        validate_size(raw)
        pdf_list.append(raw)

    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_merge_pdfs, pdf_list
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"PDF merge failed: {e}")

    return stream(data, "merged.pdf", "application/pdf")


# ── 7. Compress PDF ────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/compress-pdf")
async def compress_pdf(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    validate("compress-pdf", f.filename)
    raw = await f.read()
    validate_size(raw)
    original_size = len(raw)

    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_compress_pdf, raw
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"PDF compression failed: {e}")

    new_size = len(data)
    savings = round((1 - new_size / original_size) * 100, 1) if original_size else 0

    return StreamingResponse(
        io.BytesIO(data),
        media_type="application/pdf",
        headers={
            "Content-Disposition": 'attachment; filename="compressed.pdf"',
            "X-Original-Size":     str(original_size),
            "X-Compressed-Size":   str(new_size),
            "X-Size-Reduction":    f"{savings}%",
        },
    )

# ── 8. Excel → PDF ─────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/excel-to-pdf")
async def excel_to_pdf(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    ext = Path(f.filename or "").suffix.lower()
    if ext not in {".xls", ".xlsx"}:
        raise HTTPException(422, "Please upload an .xls or .xlsx file.")
    raw = await f.read()
    validate_size(raw)
    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_excel_to_pdf, raw, f.filename
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Excel to PDF failed: {e}")
    return stream(data, "converted.pdf", "application/pdf")


# ── 9. PPT → PDF ───────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/ppt-to-pdf")
async def ppt_to_pdf(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    ext = Path(f.filename or "").suffix.lower()
    if ext not in {".ppt", ".pptx"}:
        raise HTTPException(422, "Please upload a .ppt or .pptx file.")
    raw = await f.read()
    validate_size(raw)
    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_ppt_to_pdf, raw
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"PPT to PDF failed: {e}")
    return stream(data, "converted.pdf", "application/pdf")


# ── 10. PDF → PNG ──────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/pdf-to-png")
async def pdf_to_png(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    if Path(f.filename or "").suffix.lower() != ".pdf":
        raise HTTPException(422, "Please upload a .pdf file.")
    raw = await f.read()
    validate_size(raw)
    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_pdf_to_png, raw
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"PDF to PNG failed: {e}")
    return stream(data, "pages.zip", "application/zip")


# ── 11. Compress Image ─────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/compress-image")
async def compress_image(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    ext = Path(f.filename or "").suffix.lower()
    if ext not in {".jpg", ".jpeg", ".png"}:
        raise HTTPException(422, "Please upload a JPG or PNG file.")
    raw = await f.read()
    validate_size(raw)
    original_size = len(raw)
    try:
        data, out_name = await asyncio.get_event_loop().run_in_executor(
            None, do_compress_image, raw, f.filename or "image.jpg"
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Image compression failed: {e}")
    savings = round((1 - len(data) / original_size) * 100, 1)
    return StreamingResponse(
        io.BytesIO(data),
        media_type="image/jpeg" if out_name.endswith(".jpg") else "image/png",
        headers={
            "Content-Disposition": f'attachment; filename="{out_name}"',
            "X-Original-Size":    str(original_size),
            "X-Compressed-Size":  str(len(data)),
            "X-Size-Reduction":   f"{savings}%",
        },
    )


# ── 12. Unlock PDF ─────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/unlock-pdf")
async def unlock_pdf(
    request: Request,
    files: List[UploadFile] = File(...),
    password: str = Form(default=""),
):
    f = files[0]
    if Path(f.filename or "").suffix.lower() != ".pdf":
        raise HTTPException(422, "Please upload a .pdf file.")
    raw = await f.read()
    validate_size(raw)
    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_unlock_pdf, raw, password
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"PDF unlock failed: {e}")
    return stream(data, "unlocked.pdf", "application/pdf")


# ── 13. JPG → PNG ──────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/jpg-to-png")
async def jpg_to_png(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    ext = Path(f.filename or "").suffix.lower()
    if ext not in {".jpg", ".jpeg"}:
        raise HTTPException(422, "Please upload a JPG or JPEG file.")
    raw = await f.read()
    validate_size(raw)
    try:
        data = await asyncio.get_event_loop().run_in_executor(None, do_jpg_to_png, raw)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"JPG to PNG failed: {e}")
    name = Path(f.filename or "image").stem + ".png"
    return stream(data, name, "image/png")


# ── 14. PNG → JPG ──────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/png-to-jpg")
async def png_to_jpg(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    if Path(f.filename or "").suffix.lower() != ".png":
        raise HTTPException(422, "Please upload a PNG file.")
    raw = await f.read()
    validate_size(raw)
    try:
        data = await asyncio.get_event_loop().run_in_executor(None, do_png_to_jpg, raw)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"PNG to JPG failed: {e}")
    name = Path(f.filename or "image").stem + ".jpg"
    return stream(data, name, "image/jpeg")


# ── 15. HTML → PDF ─────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/html-to-pdf")
async def html_to_pdf(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    ext = Path(f.filename or "").suffix.lower()
    if ext not in {".html", ".htm"}:
        raise HTTPException(422, "Please upload an HTML or HTM file.")
    raw = await f.read()
    validate_size(raw)
    try:
        data = await asyncio.get_event_loop().run_in_executor(None, do_html_to_pdf, raw)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"HTML to PDF failed: {e}")
    return stream(data, "converted.pdf", "application/pdf")


# ── 16. Watermark PDF ──────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/watermark-pdf")
async def watermark_pdf(
    request: Request,
    files: List[UploadFile] = File(...),
    watermark_text: str = Form(default="CONFIDENTIAL"),
):
    f = files[0]
    if Path(f.filename or "").suffix.lower() != ".pdf":
        raise HTTPException(422, "Please upload a PDF file.")
    if not watermark_text.strip():
        raise HTTPException(422, "Watermark text cannot be empty.")
    raw = await f.read()
    validate_size(raw)
    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_watermark_pdf, raw, watermark_text.strip()
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Watermark failed: {e}")
    return stream(data, "watermarked.pdf", "application/pdf")


# ── 17. PDF → HTML ─────────────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/pdf-to-html")
async def pdf_to_html(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    if Path(f.filename or "").suffix.lower() != ".pdf":
        raise HTTPException(422, "Please upload a PDF file.")
    raw = await f.read()
    validate_size(raw)
    try:
        data = await asyncio.get_event_loop().run_in_executor(None, do_pdf_to_html, raw)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"PDF to HTML failed: {e}")
    name = Path(f.filename or "document").stem + ".html"
    return stream(data, name, "text/html")


# ── 18. Remove Background ──────────────────────────────────────────────────────
@limiter.limit(RATE_LIMIT)
@app.post("/convert/remove-background")
async def remove_background(request: Request, files: List[UploadFile] = File(...)):
    f = files[0]
    ext = Path(f.filename or "").suffix.lower()
    if ext not in {".jpg", ".jpeg", ".png"}:
        raise HTTPException(422, "Please upload a JPG or PNG file.")
    raw = await f.read()
    validate_size(raw)
    try:
        data = await asyncio.get_event_loop().run_in_executor(
            None, do_remove_background, raw
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Background removal failed: {e}")
    name = Path(f.filename or "image").stem + "_nobg.png"
    return stream(data, name, "image/png")